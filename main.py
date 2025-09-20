import json
import os
from groq import Groq
import PyPDF2
import io
from typing import List, Dict, Any
import re
from PIL import Image
import pytesseract
from docx import Document

class DRIForesightProcessor:
    def __init__(self, groq_api_key: str):
        """Initialize the DRI Foresight processor with Groq API."""
        self.client = Groq(api_key=groq_api_key)
        self.model = "meta-llama/llama-4-scout-17b-16e-instruct"  # Using available model
        
    def extract_text_from_pdf(self, pdf_file) -> str:
        """Extract text content from uploaded PDF file."""
        try:
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(pdf_file.read()))
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            return f"Error extracting PDF text: {str(e)}"

    def extract_text_from_file(self, file) -> str:
        """Extract text content from uploaded file (supports multiple formats)."""
        try:
            file_extension = file.name.split('.')[-1].lower()
            
            if file_extension == 'pdf':
                return self.extract_text_from_pdf(file)
            elif file_extension in ['txt']:
                return file.read().decode('utf-8')
            elif file_extension in ['csv']:
                # Prefer pandas if available; fallback to Python csv if not
                try:
                    import pandas as pd  # type: ignore
                    df = pd.read_csv(file)
                    return df.to_string()
                except Exception:
                    try:
                        file.seek(0)
                        import csv as _csv
                        decoded_lines = file.read().decode('utf-8', 'ignore').splitlines()
                        reader = _csv.reader(decoded_lines)
                        rows = list(reader)
                        return "\n".join([", ".join(row) for row in rows])
                    except Exception as csv_err:
                        return f"Could not read CSV file: {csv_err}"
            #newly added from this 
            elif file_extension in ['docx']:
                doc = Document(file)
                text = ""
                for paragraph in doc.paragraphs:
                    text += paragraph.text + "\n"
                return text.strip()

            elif file_extension in ['doc']:
                # For .doc files, you might need python-docx2txt
                try:
                    import docx2txt
                    return docx2txt.process(file)
                except ImportError:
                    return "docx2txt library required for .doc files"

            elif file_extension in ['pptx']:
                from pptx import Presentation
                prs = Presentation(file)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text.strip()

            elif file_extension in ['ppt']:
                # For .ppt files, you might need additional libraries like python-pptx or comtypes
                return "PowerPoint .ppt format requires additional processing"

            elif file_extension in ['jpg', 'jpeg', 'png', 'bmp', 'gif']:
                # OCR for images
                image = Image.open(file)
                text = pytesseract.image_to_string(image)
                return text.strip()
            #to this 
            else:
                # For other formats, try to read as text
                try:
                    return file.read().decode('utf-8')
                except:
                    return f"Could not extract text from {file.name}"
        except Exception as e:
            return f"Error extracting text from {file.name}: {str(e)}"

    def generate_domain_map(self, domain: str, document_text: str, project_name: str) -> Dict[str, Any]:
        """Generate domain map based on the selected domain and document content."""
        
        # Check if we have substantial document content
        has_document_content = document_text and len(document_text.strip()) > 100
        
        if has_document_content:
            # Prioritize document content when available
            prompt = f"""
            You are an expert in foresight analysis and strategic planning. Based on the project "{project_name}", 
            analyze the following document content to create a comprehensive domain map. The selected domain focus "{domain}" 
            should be used as context, but the domain map should primarily reflect the content and themes found in the uploaded document.
            Please generate exactly 5-7 sub-domains to ensure comprehensive coverage of the domain.
            

            Document Content:
            {document_text[:3000]}

            Please analyze the document content thoroughly and generate a detailed domain map with 5-7 sub-domains that reflects the actual themes and topics discussed in the document.

            For each sub-domain:
            1. Provide a detailed description (2-3 sentences) that captures the specific aspects discussed in the document
            2. Identify 8-12 specific issue areas that are either mentioned in the document or are relevant challenges within that sub-domain
            3. Base the descriptions and issue areas on the actual content rather than generic knowledge

            Format your response as a JSON object with this structure:
            {{
                "central_domain": "Main focus area based on document content",
                "description": "Detailed description derived from the document content (2-3 sentences)",
                "sub_domains": [
                    {{
                        "name": "Sub-domain name from document themes",
                        "description": "Comprehensive description based on specific document content (2-3 sentences explaining what the document reveals about this area)",
                        "relevance": "High/Medium/Low",
                        "issue_areas": [
                            "Specific issue area 1 from document context",
                            "Specific issue area 2 from document context",
                            "Specific issue area 3 from document context",
                            "Specific issue area 4 from document context",
                            "Specific issue area 5 from document context",
                            "Specific issue area 6 from document context",
                            "Specific issue area 7 from document context",
                            "Specific issue area 8 from document context"
                        ]
                    }}
                ]
            }}

            Focus on what the document actually discusses in detail rather than providing generic overviews.
            """

        else:
            # Fall back to domain-based generation when no substantial document content
            prompt = f"""
            You are an expert in foresight analysis and strategic planning. Based on the project "{project_name}" focusing on the domain "{domain}", 
            create a comprehensive and detailed domain map for this specific focus area.
            Please generate exactly 5-7 sub-domains to ensure comprehensive coverage of the domain.

            Generate detailed analysis with:
            1. 5-7 comprehensive sub-domains with detailed descriptions (2-3 sentences each)
            2. Comprehensive sub-domain descriptions that explain the specific aspects and importance of each area
            3. 8-12 specific issue areas per sub-domain that represent real challenges, opportunities, or focus points

            Format your response as a JSON object with this structure:
            {{
                "central_domain": "Main focus area title",
                "description": "Detailed description of the central domain (2-3 sentences)",
                "sub_domains": [
                    {{
                        "name": "Sub-domain name",
                        "description": "Comprehensive description explaining the specific aspects, challenges, and importance of this sub-domain (2-3 sentences)",
                        "relevance": "High/Medium/Low",
                        "issue_areas": [
                            "Specific issue area 1 with clear focus",
                            "Specific issue area 2 with clear focus",
                            "Specific issue area 3 with clear focus",
                            "Specific issue area 4 with clear focus",
                            "Specific issue area 5 with clear focus",
                            "Specific issue area 6 with clear focus",
                            "Specific issue area 7 with clear focus",
                            "Specific issue area 8 with clear focus"
                        ]
                    }}
                ]
            }}

            Provide detailed, actionable descriptions and specific issue areas rather than generic summaries.
            """
                
        try:
            chat_completion = self.client.chat.completions.create(
                messages=[
                    {
                        "role": "system", 
                        "content": "You are an expert foresight analyst specializing in domain mapping and strategic analysis. Always respond with valid JSON format."
                    },
                    {"role": "user", "content": prompt}
                ],
                model=self.model,
                max_tokens=1500,
                temperature=0.7
            )
            
            response_text = chat_completion.choices[0].message.content
            return self._parse_json_response(response_text)
                
        except Exception as e:
            return {"error": f"Failed to generate domain map: {str(e)}"}
   
    # UPDATE 1: In generate_signals method - around line 120
    def generate_signals(self, domain: str, document_text: str) -> Dict[str, List[Dict]]:
        """Generate strong and weak signals based on document analysis including interview insights."""
        
        # UPDATED: Enhanced prompt to better handle comprehensive document context
        prompt = f"""
        As a foresight expert analyzing the domain "{domain}", examine the following comprehensive content which includes:
        - Domain mapping documents and project materials
        - Interview transcripts and stakeholder insights  
        - External signals and trend data
        - Any additional research materials
        
        COMPREHENSIVE ANALYSIS TASK:
        1. STRONG SIGNALS: Clear, evident trends or changes that are already happening
        2. WEAK SIGNALS: Early indicators of potential future changes that might be emerging

        COMPREHENSIVE CONTENT (ALL SOURCES COMBINED):
        {document_text[:12000]}  # Increased limit to capture more content

        ANALYSIS INSTRUCTIONS:
        - Synthesize insights across ALL uploaded content types
        - Pay special attention to interview insights for stakeholder perspectives
        - Look for patterns and convergence across different data sources
        - Include signals that emerge from cross-referencing different document types
        - Clearly indicate source context in descriptions

        Please provide 5-7 strong signals and 5-7 weak signals in the following JSON format:
        {{
            "strong_signals": [
                {{
                    "title": "Signal title",
                    "description": "Detailed description synthesizing multiple sources where relevant",
                    "source": "Primary source type (documents/interviews/signals/cross-source pattern)",
                    "impact": "Potential impact description",
                    "evidence_strength": "Evidence level from uploaded materials"
                }}
            ],
            "weak_signals": [
                {{
                    "title": "Signal title", 
                    "description": "Detailed description synthesizing multiple sources where relevant",
                    "source": "Primary source type (documents/interviews/signals/cross-source pattern)",
                    "potential": "Future potential or implications",
                    "evidence_strength": "Evidence level from uploaded materials"
                }}
            ]
        }}

        CRITICAL REQUIREMENTS:
        - Every signal must be grounded in the provided content
        - Prioritize signals that appear across multiple source types
        - Include stakeholder perspectives from interviews where available
        - Focus on domain-specific insights
        """
        
        try:
            chat_completion = self.client.chat.completions.create(
                messages=[
                    {
                        "role": "system", 
                        "content": """You are an expert in comprehensive signal detection and trend analysis for strategic foresight. 
                        You excel at synthesizing insights from multiple data sources including documents, interviews, 
                        external signals, and research materials. You always identify patterns across different source types.
                        Always respond with valid JSON format."""
                    },
                    {"role": "user", "content": prompt}
                ],
                model=self.model,
                max_tokens=2500,  # Increased for comprehensive analysis
                temperature=0.8
            )
            
            response_text = chat_completion.choices[0].message.content
            return self._parse_json_response(response_text)
                
        except Exception as e:
            return {"error": f"Failed to generate signals: {str(e)}"}


    # UPDATE 2: In generate_steepv_analysis method - around line 180
    def generate_steepv_analysis(self, domain: str, signals_data: Dict, document_text: str) -> Dict[str, List[str]]:
        """Generate comprehensive STEEPV analysis based on signals, domain, and all available context."""
        
        # UPDATED: Better signal extraction and handling
        strong_signals = signals_data.get('strong_signals', [])
        weak_signals = signals_data.get('weak_signals', [])
        
        # Enhanced signal processing
        if isinstance(signals_data, dict) and 'raw_response' in signals_data:
            try:
                import json
                import re
                raw_response = signals_data['raw_response']
                json_match = re.search(r'```json\s*(\{.*?\})\s*```', raw_response, re.DOTALL)
                if json_match:
                    parsed_signals = json.loads(json_match.group(1))
                    strong_signals = parsed_signals.get('strong_signals', [])
                    weak_signals = parsed_signals.get('weak_signals', [])
            except:
                pass
        
        # Create comprehensive signal descriptions
        signal_descriptions = []
        for signal in strong_signals:
            if isinstance(signal, dict):
                title = signal.get('title', 'Unknown')
                desc = signal.get('description', '')
                source = signal.get('source', '')
                signal_descriptions.append(f"STRONG: {title} - {desc} (Source: {source})")
            else:
                signal_descriptions.append(f"STRONG: {signal}")
        
        for signal in weak_signals:
            if isinstance(signal, dict):
                title = signal.get('title', 'Unknown')
                desc = signal.get('description', '')
                source = signal.get('source', '')
                signal_descriptions.append(f"WEAK: {title} - {desc} (Source: {source})")
            else:
                signal_descriptions.append(f"WEAK: {signal}")
        
        # UPDATED: Enhanced prompt with comprehensive context integration
        prompt = f"""
        You are conducting a comprehensive STEEPV analysis for the domain: "{domain}"

        ANALYSIS CONTEXT:
        - Domain Focus: {domain}
        - Identified Signals: {len(signal_descriptions)} signals from multiple sources
        - Comprehensive Context: Documents, interviews, external signals, research materials

        SIGNALS TO CATEGORIZE:
        {chr(10).join(signal_descriptions[:20])}  # Increased signal limit

        FULL INTEGRATED CONTEXT (All uploaded materials combined):
        {document_text[:12000] if document_text else "No additional context provided"}

        STEEPV ANALYSIS TASK:
        Analyze ALL available information and provide 4-6 specific factors for EACH STEEPV category.
        Synthesize insights from:
        - Domain mapping documents
        - Stakeholder interviews and perspectives  
        - External signals and trends
        - Research materials and data
        - Cross-source patterns and themes

        STEEPV FRAMEWORK (Enhanced Definitions):
        - Social: Demographics, cultural shifts, social movements, community behaviors, stakeholder perspectives, social challenges from interviews
        - Technological: Digital innovations, emerging technologies, automation, AI, technical barriers and opportunities from all sources
        - Economic: Market conditions, funding landscapes, costs, financial challenges/opportunities, economic trends from interviews and documents
        - Environmental: Climate factors, sustainability requirements, environmental concerns from stakeholder input and research
        - Political: Government policies, regulatory environment, political factors, governance challenges from comprehensive analysis
        - Values: Ethical frameworks, cultural values, stakeholder beliefs, value systems from interviews and cultural analysis

        COMPREHENSIVE ANALYSIS REQUIREMENTS:
        1. Each category must contain 4-6 specific, actionable factors
        2. Factors must be grounded in the provided materials (documents + interviews + signals)
        3. Prioritize factors that appear across multiple source types
        4. Include stakeholder perspectives from interviews where relevant
        5. Make factors specific to the "{domain}" domain context
        6. Ensure comprehensive coverage - NO category left empty
        7. Focus on factors that will impact future scenario development

        REQUIRED JSON FORMAT:
        {{
            "Social": ["factor1 (source context)", "factor2 (source context)", "factor3 (source context)", "factor4 (source context)"],
            "Technological": ["factor1 (source context)", "factor2 (source context)", "factor3 (source context)", "factor4 (source context)"],
            "Economic": ["factor1 (source context)", "factor2 (source context)", "factor3 (source context)", "factor4 (source context)"],
            "Environmental": ["factor1 (source context)", "factor2 (source context)", "factor3 (source context)", "factor4 (source context)"],
            "Political": ["factor1 (source context)", "factor2 (source context)", "factor3 (source context)", "factor4 (source context)"],
            "Values": ["factor1 (source context)", "factor2 (source context)", "factor3 (source context)", "factor4 (source context)"]
        }}

        CRITICAL: Every category must be populated with content-grounded factors from the comprehensive materials provided.
        """
        
        try:
            chat_completion = self.client.chat.completions.create(
                messages=[
                    {
                        "role": "system",
                        "content": """You are a senior strategic foresight analyst specializing in comprehensive STEEPV methodology. 
                        You excel at synthesizing multiple data sources (documents, interviews, signals, research) into structured analysis.
                        Your expertise is in ensuring complete coverage across all STEEPV dimensions using integrated evidence.
                        Always respond with valid, complete JSON containing well-grounded factors for all categories."""
                    },
                    {"role": "user", "content": prompt}
                ],
                model=self.model,
                max_tokens=3000,  # Increased for comprehensive analysis
                temperature=0.6
            )
            
            response_text = chat_completion.choices[0].message.content
            parsed_result = self._parse_json_response(response_text)
            
            # Enhanced validation with domain-specific fallbacks
            steepv_categories = ["Social", "Technological", "Economic", "Environmental", "Political", "Values"]
            
            for category in steepv_categories:
                if category not in parsed_result or not parsed_result[category]:
                    # Domain-specific fallbacks based on comprehensive analysis
                    parsed_result[category] = [
                        f"{category} factors identified in {domain} domain analysis",
                        f"{category} implications from stakeholder interviews", 
                        f"{category} trends affecting {domain} development",
                        f"{category} considerations from uploaded materials"
                    ]
            
            return parsed_result
                    
        except Exception as e:
            # Enhanced fallback with domain context
            return {
                "Social": [f"Social dynamics in {domain} from interviews", "Community engagement patterns from analysis", "Cultural factors from comprehensive review"],
                "Technological": [f"Technology adoption in {domain}", "Digital transformation patterns", "Innovation barriers from stakeholder input"],
                "Economic": [f"Economic conditions affecting {domain}", "Funding challenges from interviews", "Cost factors from document analysis"],
                "Environmental": [f"Environmental considerations in {domain}", "Sustainability requirements from research", "Climate impacts from comprehensive analysis"],
                "Political": [f"Policy environment for {domain}", "Regulatory factors from documents", "Governance challenges from interviews"],
                "Values": [f"Value systems in {domain} context", "Ethical frameworks from analysis", "Cultural alignment from stakeholder input"]
            }

    def _parse_json_response(self, response_text: str) -> Dict:
        """Enhanced JSON parsing with better error handling."""
        import json
        import re
        
        # def clean_json_string(json_str):
        #     """Clean up common JSON formatting issues"""
        #     # Fix newlines immediately after opening quotes
        #     json_str = re.sub(r':\s*"\s*\n\s*', ': "', json_str)
        #     # Fix multiple newlines within strings
        #     json_str = re.sub(r'\n\s*\n', '\\n\\n', json_str)
        #     # Fix single newlines within strings (but preserve paragraph structure)
        #     json_str = re.sub(r'(?<!\\)(?<!\\n)\n(?!\s*[}\]",])', ' ', json_str)

        #     # Remove markdown fences
        #     json_str = re.sub(r"```(?:json)?", "", json_str)
        #     json_str = json_str.replace("```", "")

        #     # Remove trailing commas before } or ]
        #     json_str = re.sub(r",\s*([\]}])", r"\1", json_str)

        #     # Normalize multiple spaces/newlines
        #     json_str = re.sub(r"\s+\n", " ", json_str)
        #     json_str = re.sub(r"\n+", " ", json_str)

        #     # return json_str
        #     return json_str.strip()

        def clean_json_string(json_str):
            """Clean up common JSON formatting issues"""
            # Fix newlines immediately after opening quotes
            json_str = re.sub(r':\s*"\s*\n\s*', ': "', json_str)
            # Fix multiple newlines within strings
            json_str = re.sub(r'\n\s*\n', '\\n\\n', json_str)
            # Fix single newlines within strings
            json_str = re.sub(r'(?<!\\)(?<!\\n)\n(?!\s*[}\]",])', ' ', json_str)

            # Remove markdown fences
            json_str = re.sub(r"```(?:json)?", "", json_str)
            json_str = json_str.replace("```", "")

            # Fix missing commas between objects
            json_str = re.sub(r'\}\s*\{', '}, {', json_str)

            # Remove dangling/trailing commas
            json_str = re.sub(r",\s*([\]}])", r"\1", json_str)
            json_str = re.sub(r',\s*,+', ',', json_str)

            # Escape unescaped quotes inside values
            json_str = re.sub(
                r'(?<=:\s")([^"]*?)"(?=\s*[,}])',
                lambda m: m.group(1).replace('"', '\\"'),
                json_str
            )

            # Normalize multiple spaces/newlines
            json_str = re.sub(r"\s+\n", " ", json_str)
            json_str = re.sub(r"\n+", " ", json_str)

            return json_str.strip()

        
        try:
            # First try direct JSON parsing
            return json.loads(response_text)
        except:
            try:
                # Look for JSON in code blocks - more flexible pattern
                json_match = re.search(r'```(?:json)?\s*\n?(.*?)\n?\s*```', response_text, re.DOTALL)
                if json_match:
                    extracted_json = json_match.group(1).strip()
                    # Try to clean up formatting issues
                    cleaned_json = clean_json_string(extracted_json)
                    try:
                        return json.loads(cleaned_json)
                    except:
                        # If cleaning didn't work, try original
                        return json.loads(extracted_json)
                
                # Look for JSON-like structure without code blocks
                json_match = re.search(r'\{.*\}', response_text, re.DOTALL)
                if json_match:
                    extracted_json = json_match.group(0)
                    cleaned_json = clean_json_string(extracted_json)
                    try:
                        return json.loads(cleaned_json)
                    except:
                        return json.loads(extracted_json)
                    
            except Exception as e:
                print(f"JSON parsing error: {e}")
                print(f"Problematic JSON snippet: {response_text[:500]}...")
                
            # Return empty dict if all parsing fails
            return {}

    def generate_ai_suggestions(self, domain: str, signals_data: Dict) -> List[Dict]:
        """Generate AI-powered suggestions for additional signals to consider."""
        
        existing_signals = []
        for signal in signals_data.get('strong_signals', []) + signals_data.get('weak_signals', []):
            existing_signals.append(signal.get('title', ''))
        
        prompt = f"""
        Given the domain "{domain}" and the following existing signals, suggest 3-5 additional signals 
        that should be monitored for comprehensive foresight analysis.

        Existing Signals:
        {chr(10).join(existing_signals)}

        Provide suggestions for signals that:
        1. Are not already covered
        2. Are relevant to the domain
        3. Could significantly impact future scenarios
        4. Come from different perspectives or sectors

        Format as JSON:
        {{
            "suggestions": [
                {{
                    "title": "Suggested signal title",
                    "description": "Why this signal is important to monitor",
                    "category": "Strong/Weak",
                    "rationale": "Why this wasn't covered in existing signals"
                }}
            ]
        }}
        """
        
        try:
            chat_completion = self.client.chat.completions.create(
                messages=[
                    {
                        "role": "system",
                        "content": "You are an expert in signal detection and strategic foresight analysis. Always respond with valid JSON format."
                    },
                    {"role": "user", "content": prompt}
                ],
                model=self.model,
                max_tokens=1000,
                temperature=0.8
            )
            
            response_text = chat_completion.choices[0].message.content
            parsed_response = self._parse_json_response(response_text)
            
            if 'error' in parsed_response:
                return [{"error": parsed_response['error'], "raw_response": parsed_response.get('raw_response', '')}]
            
            return parsed_response.get('suggestions', [])
                
        except Exception as e:
            return [{"error": f"Failed to generate suggestions: {str(e)}"}]

    # UPDATE 3: In generate_futures_triangle method - around line 320
    def generate_futures_triangle(self, domain: str, signals_data: Dict, steepv_data: Dict, interview_context: str = "") -> Dict[str, Any]:
        """Generate comprehensive Futures Triangle analysis based on all available data sources."""
        
        # Extract signals for context
        strong_signals = signals_data.get('strong_signals', [])
        weak_signals = signals_data.get('weak_signals', [])
        
        # Enhanced signal formatting with source context
        strong_signals_text = "\n".join([
            f"- {signal.get('title', '')}: {signal.get('description', '')} (Source: {signal.get('source', 'Analysis')})" 
            for signal in strong_signals
        ])
        weak_signals_text = "\n".join([
            f"- {signal.get('title', '')}: {signal.get('description', '')} (Source: {signal.get('source', 'Analysis')})" 
            for signal in weak_signals
        ])
        
        # Enhanced STEEPV formatting
        steepv_text = ""
        for category, factors in steepv_data.items():
            if factors:
                steepv_text += f"\n{category}: {', '.join(factors[:4])}"
        
        # UPDATED: Comprehensive interview and document integration
        comprehensive_context = ""
        if interview_context:
            comprehensive_context = f"""
            
            COMPREHENSIVE DOCUMENT CONTEXT (All Sources):
            {interview_context[:8000]}  # Increased limit for full context
            """
        
        # UPDATED: Enhanced prompt for comprehensive analysis with Key Dynamics
        prompt = f"""
        As a strategic foresight analyst, create a comprehensive Futures Triangle analysis for the domain "{domain}".

        INTEGRATED ANALYSIS BASE:

        STRONG SIGNALS (from comprehensive analysis):
        {strong_signals_text}

        WEAK SIGNALS (from comprehensive analysis):
        {weak_signals_text}

        STEEPV ANALYSIS SUMMARY:
        {steepv_text}
        {comprehensive_context}

        FUTURES TRIANGLE METHODOLOGY:
        Create a comprehensive analysis integrating ALL uploaded materials (domain documents, interviews, signals, research) into the three temporal forces:

        1. PULL OF THE FUTURE (Emerging Issues & Aspirations):
        - Weak Signals: Early indicators of possible change (experiments, anomalies, fringe innovations)
        - Emerging Issues: New challenges or opportunities just becoming visible
        - Visions & Aspirations: Images of preferred futures and goals pulling society forward

        2. PUSH OF THE PRESENT (Current Momentum & Drivers):
        - Current Trends: Observable patterns of change with clear direction
        - Strong Drivers: Active forces creating pressure for change

        3. WEIGHT OF HISTORY (Historical Constraints & Values):
        - Barriers & Inertia: Structures and systems resisting change + tendency to continue current patterns (e.g., laws, infrastructure gaps, financial limits, institutional routines, cultural habits, organizational momentum)
        - Values to Preserve: Elements worth preserving through change (e.g., democratic principles, cultural heritage)

        4. KEY DYNAMICS & STRATEGIC INSIGHTS:
        - Primary Tensions: Main conflicts between the three forces
        - Alignment Opportunities: Where forces work together effectively
        - Critical Uncertainties: What remains unknown or unpredictable

        COMPREHENSIVE INTEGRATION REQUIREMENTS:
        - Synthesize insights from domain documents, stakeholder interviews, external signals, and research materials
        - Ensure each force reflects evidence from multiple source types
        - Include stakeholder perspectives prominently in future visions
        - Ground all factors in the comprehensive materials provided
        - Focus on domain-specific temporal dynamics

        FORMAT YOUR RESPONSE AS JSON:
        {{
            "pull_of_future": {{
                "weak_signals": [
                    "weak signal 1 (source context)",
                    "weak signal 2 (source context)",
                    "weak signal 3 (source context)"
                ],
                "emerging_issues": [
                    "emerging issue 1 (source context)",
                    "emerging issue 2 (source context)",
                    "emerging issue 3 (source context)"
                ],
                "visions_and_aspirations": [
                    "vision/aspiration 1 (source context)",
                    "vision/aspiration 2 (source context)",
                    "vision/aspiration 3 (source context)",
                    "vision/aspiration 4 (source context)"
                ]

            }},
            "push_of_present": {{
                "current_trends": [
                    "current trend 1 (source context)",
                    "current trend 2 (source context)",
                    "current trend 3 (source context)"
                ],
                "strong_drivers": [
                    "strong driver 1 (source context)",
                    "strong driver 2 (source context)",
                    "strong driver 3 (source context)"
                ]
            }},
            "weight_of_history": {{
                "barriers_and_inertia": [
                    "barrier/inertia 1 (source context)",
                    "barrier/inertia 2 (source context)",
                    "barrier/inertia 3 (source context)",
                    "barrier/inertia 4 (source context)"
                ],
                "values_to_preserve": [
                    "value to preserve 1 (source context)",
                    "value to preserve 2 (source context)",
                    "value to preserve 3 (source context)"
                ]
            }},
            "key_dynamics": {{
                "primary_tensions": [
                    "primary tension 1 (source context)",
                    "primary tension 2 (source context)",
                    "primary tension 3 (source context)"
                ],
                "alignment_opportunities": [
                    "alignment opportunity 1 (source context)",
                    "alignment opportunity 2 (source context)",
                    "alignment opportunity 3 (source context)"
                ],
                "critical_uncertainties": [
                    "critical uncertainty 1 (source context)",
                    "critical uncertainty 2 (source context)",
                    "critical uncertainty 3 (source context)"
                ]
            }}
        }}

        Ensure each subcategory has 3-4 specific factors grounded in the comprehensive {domain} analysis.
        """
        
        try:
            chat_completion = self.client.chat.completions.create(
                messages=[
                    {
                        "role": "system",
                        "content": """You are an expert in comprehensive futures studies and the Futures Triangle methodology. 
                        You excel at integrating multiple data sources (documents, interviews, signals, research) into temporal analysis.
                        Your expertise is in synthesizing diverse materials into coherent past-present-future dynamics with strategic insights.
                        Always respond with valid JSON format grounded in provided evidence."""
                    },
                    {"role": "user", "content": prompt}
                ],
                model=self.model,
                max_tokens=3000,  # Increased for comprehensive output including key dynamics
                temperature=0.7
            )
            
            response_text = chat_completion.choices[0].message.content
            return self._parse_json_response(response_text)
                
        except Exception as e:
            return {"error": f"Failed to generate comprehensive futures triangle: {str(e)}"}

    # UPDATE 4: Add new method for comprehensive text extraction
    def extract_comprehensive_text(self, files_dict: Dict) -> str:
        """Extract and combine text from all uploaded file types for comprehensive analysis."""
        all_text_content = []
        
        # Process domain map documents
        if files_dict.get('documents'):
            all_text_content.append("=== DOMAIN MAPPING DOCUMENTS ===")
            for file in files_dict['documents']:
                content = self.extract_text_from_file(file)
                all_text_content.append(f"Document: {file.name}")
                all_text_content.append(content)
                all_text_content.append("---")
        
        # Process interview data
        if files_dict.get('interviews'):
            all_text_content.append("=== INTERVIEW DATA & STAKEHOLDER INSIGHTS ===")
            for file in files_dict['interviews']:
                content = self.extract_text_from_file(file)
                all_text_content.append(f"Interview Source: {file.name}")
                all_text_content.append(content)
                all_text_content.append("---")
        
        # Process external signals
        if files_dict.get('signals'):
            all_text_content.append("=== EXTERNAL SIGNALS & TREND DATA ===")
            for file in files_dict['signals']:
                content = self.extract_text_from_file(file)
                all_text_content.append(f"Signal Source: {file.name}")
                all_text_content.append(content)
                all_text_content.append("---")
        
        # Process domain map file separately if exists
        if files_dict.get('domain_map'):
            all_text_content.append("=== DOMAIN MAP REFERENCE ===")
            content = self.extract_text_from_file(files_dict['domain_map'])
            all_text_content.append(f"Domain Map: {files_dict['domain_map'].name}")
            all_text_content.append(content)
            all_text_content.append("---")
        
        return "\n".join(all_text_content)
    
    def analyze_interview_data(self, domain: str, interview_text: str) -> Dict[str, Any]:
        """Analyze interview data to extract challenges, opportunities, and visions."""
        
        prompt = f"""
        As an expert analyst, analyze the following interview data for the domain "{domain}".
        
        Interview Content:
        {interview_text[:4000]}  # Limit content to avoid token limits
        
        Extract and categorize the key insights into:
        1. Top Challenges - main obstacles, problems, or difficulties mentioned
        2. Key Opportunities - opportunities, potential solutions, or positive developments
        3. Future Visions - aspirations, goals, or desired future states mentioned
        
        Format your response as JSON:
        {{
            "challenges": [
                "challenge 1",
                "challenge 2",
                "challenge 3",
                "challenge 4"
            ],
            "opportunities": [
                "opportunity 1",
                "opportunity 2", 
                "opportunity 3",
                "opportunity 4"
            ],
            "visions": [
                "vision 1",
                "vision 2",
                "vision 3",
                "vision 4"
            ]
        }}
        
        Focus on the most significant and frequently mentioned themes.
        """
        
        try:
            chat_completion = self.client.chat.completions.create(
                messages=[
                    {
                        "role": "system",
                        "content": "You are an expert qualitative data analyst specializing in interview analysis and thematic extraction. Always respond with valid JSON format."
                    },
                    {"role": "user", "content": prompt}
                ],
                model=self.model,
                max_tokens=1500,
                temperature=0.6
            )
            
            response_text = chat_completion.choices[0].message.content
            return self._parse_json_response(response_text)
                
        except Exception as e:
            return {"error": f"Failed to analyze interview data: {str(e)}"}

    def generate_futures_triangle_2_0(self, domain: str, phase1_data: Dict, phase2_data: Dict, comprehensive_context: str = "") -> Dict[str, Any]:
            """Generate enhanced Futures Triangle 2.0 analysis for Phase 3 scenario planning."""
            
            # Extract Phase 2 data
            signals_data = phase2_data.get('signals_data', {})
            steepv_data = phase2_data.get('steepv_data', {})
            basic_triangle = phase2_data.get('futures_triangle_data', {})
            
            # Format signals context
            strong_signals = signals_data.get('strong_signals', [])
            weak_signals = signals_data.get('weak_signals', [])
            
            signals_context = ""
            if strong_signals:
                signals_context += "STRONG SIGNALS:\n" + "\n".join([
                    f"- {s.get('title', '')}: {s.get('description', '')}" for s in strong_signals
                ])
            if weak_signals:
                signals_context += "\n\nWEAK SIGNALS:\n" + "\n".join([
                    f"- {s.get('title', '')}: {s.get('description', '')}" for s in weak_signals
                ])
            
            # Format STEEPV context
            steepv_context = ""
            for category, factors in steepv_data.items():
                if factors:
                    steepv_context += f"\n{category.upper()}: {', '.join(factors[:4])}"
            
            # Enhanced prompt for Futures Triangle 2.0
            prompt = f"""
            As a strategic foresight expert, create a comprehensive Futures Triangle 2.0 Analysis for "{domain}" that will directly feed into scenario planning.

            CONTEXT FROM PREVIOUS PHASES:
            Project: {phase1_data.get('project_name', domain)}
            Domain Focus: {domain}
            
            PHASE 2 ANALYSIS RESULTS:
            {signals_context}
            
            STEEPV ANALYSIS:
            {steepv_context}
            
            COMPREHENSIVE DOCUMENT CONTEXT:
            {comprehensive_context[:8000]}
            
            FUTURES TRIANGLE 2.0 METHODOLOGY:
            This enhanced version extracts three key elements for scenario building:

            1. **DRIVERS** (Enhanced from Push of Present + STEEPV):
            - Major forces creating change pressure
            - Rate each by impact level (High/Medium/Low) and certainty (High/Medium/Low)
            - These will be "bent" to different archetypes in scenario planning

            2. **UNCERTAINTIES** (Critical unknowns from analysis):
            - High-impact variables that could go multiple directions
            - Key pivot points that determine scenario outcomes
            - Wild cards and game-changing possibilities

            3. **NARRATIVES** (Stories shaping the domain):
            - Dominant mental models currently operating
            - Emerging alternative narratives from weak signals
            - Competing storylines about the future

            ENHANCED TRIANGLE STRUCTURE:
            Also provide the expanded traditional triangle with Key Dynamics for strategic insights.

            FORMAT AS JSON:
            {{
                "drivers": [
                    {{
                        "id": "D1",
                        "name": "Driver name",
                        "description": "Detailed description of the driving force",
                        "category": "Technological/Economic/Social/Environmental/Political/Values",
                        "impact_level": "High/Medium/Low",
                        "certainty": "High/Medium/Low",
                        "current_trajectory": "Current direction and momentum",
                        "source_evidence": "Evidence from uploaded materials"
                    }}
                ],
                "uncertainties": [
                    {{
                        "id": "U1",
                        "name": "Uncertainty name",
                        "description": "What is uncertain and why it matters",
                        "key_variables": ["Variable 1", "Variable 2", "Variable 3"],
                        "possible_outcomes": ["Outcome A", "Outcome B", "Outcome C"],
                        "impact_on_scenarios": "How this shapes different futures",
                        "source_evidence": "Evidence from analysis"
                    }}
                ],
                "narratives": [
                    {{
                        "id": "N1",
                        "type": "Dominant/Emerging/Alternative",
                        "name": "Narrative name",
                        "description": "The story or mental model",
                        "supporting_evidence": ["Evidence 1", "Evidence 2"],
                        "influence_areas": ["Area 1", "Area 2"],
                        "alternative_versions": ["Alternative view 1", "Alternative view 2"],
                        "source_context": "Where this narrative appears in materials"
                    }}
                ],
                "enhanced_triangle": {{
                    "pull_of_future": {{
                        "weak_signals": ["signal 1", "signal 2", "signal 3"],
                        "emerging_issues": ["issue 1", "issue 2", "issue 3"],
                        "visions_aspirations": ["vision 1", "vision 2", "vision 3"]
                    }},
                    "push_of_present": {{
                        "trends": ["trend 1", "trend 2", "trend 3"],
                        "drivers": ["driver 1", "driver 2", "driver 3"]
                    }},
                    "weight_of_history": {{
                        "barriers_inertia": ["barrier 1", "barrier 2", "barrier 3"],
                        "values_to_maintain": ["value 1", "value 2", "value 3"]
                    }},
                    "key_dynamics": {{
                        "primary_tensions": ["tension 1", "tension 2", "tension 3"],
                        "alignment_opportunities": ["opportunity 1", "opportunity 2"],
                        "critical_uncertainties": ["uncertainty 1", "uncertainty 2"]
                    }}
                }},
                "strategic_insights": {{
                    "leverage_points": ["point 1", "point 2", "point 3"],
                    "signals_to_monitor": ["signal 1", "signal 2"],
                    "values_to_protect": ["value 1", "value 2"]
                }}
            }}

            CRITICAL REQUIREMENTS:
            - Extract 4-6 DRIVERS that will dominate the baseline scenario
            - Identify 3-5 UNCERTAINTIES that are pivot points for different outcomes  
            - Capture 3-4 NARRATIVES (mix of dominant and emerging) that frame stakeholder thinking
            - Ground all elements in the provided evidence from Phase 1 & 2
            - Ensure drivers/uncertainties/narratives can be "bent" to collapse/new equilibrium/transformation archetypes
            """
            
            try:
                chat_completion = self.client.chat.completions.create(
                    messages=[
                        {
                            "role": "system",
                            "content": """You are a senior strategic foresight analyst specializing in Futures Triangle 2.0 methodology 
                            for scenario planning. You excel at extracting manipulable elements (drivers, uncertainties, narratives) 
                            from comprehensive foresight analysis that can be adapted across different scenario archetypes.
                            Always respond with valid, complete JSON."""
                        },
                        {"role": "user", "content": prompt}
                    ],
                    model=self.model,
                    max_tokens=4000,
                    temperature=0.7
                )
                
                response_text = chat_completion.choices[0].message.content
                parsed_result = self._parse_json_response(response_text)

                
                # Validate required sections exist
                required_sections = ['drivers', 'uncertainties', 'narratives', 'enhanced_triangle']
                for section in required_sections:
                    if section not in parsed_result:
                        parsed_result[section] = []
                
                return parsed_result
                    
            except Exception as e:
                return {"error": f"Failed to generate Futures Triangle 2.0: {str(e)}"}

#new
    def generate_baseline_scenario(self, domain: str, triangle_2_0_data: Dict, phase1_data: Dict = None) -> Dict[str, Any]:
        """Generate baseline scenario dominated by Push of Present and key Drivers (3-4 paragraphs)."""
        
        # Extract key elements from Futures Triangle 2.0
        drivers = triangle_2_0_data.get('drivers', [])
        enhanced_triangle = triangle_2_0_data.get('enhanced_triangle', {})
        push_of_present = enhanced_triangle.get('push_of_present', {})
        
        # Format drivers context - focus on high certainty/high impact
        high_certainty_drivers = [d for d in drivers if d.get('certainty', '').lower() in ['high', 'medium']]
        drivers_context = ""
        for driver in high_certainty_drivers[:5]:  # Limit to top 5 drivers
            drivers_context += f"- {driver.get('name', '')}: {driver.get('description', '')} (Impact: {driver.get('impact_level', '')}, Certainty: {driver.get('certainty', '')})\n"
        
        # Format Push of Present context
        trends = push_of_present.get('trends', [])
        existing_drivers = push_of_present.get('drivers', [])
        push_context = ""
        if trends:
            push_context += "Current Trends: " + ", ".join(trends[:4])
        if existing_drivers:
            push_context += "\nExisting Momentum: " + ", ".join(existing_drivers[:4])
        
        # Project context
        project_name = phase1_data.get('project_name', domain) if phase1_data else domain
        
        prompt = f"""
        You are a strategic foresight expert creating a Baseline Scenario for "{domain}".

        PROJECT CONTEXT:
        Project: {project_name}
        Domain: {domain}
        
        BASELINE SCENARIO DEFINITION:
        The baseline represents the "business-as-usual" future - what happens if current momentum continues without major surprises, disruptions, or transformative changes. This is dominated by:
        1. **Push of the Present**: Current trends and momentum
        2. **Key Drivers**: High-certainty forces shaping the future
        
        PUSH OF THE PRESENT (Current Momentum):
        {push_context}
        
        KEY DRIVERS (High Certainty Forces):
        {drivers_context}
        
        BASELINE SCENARIO REQUIREMENTS:
        
        **Structure (3-4 paragraphs, 400-450 words total):**
        
        **Paragraph 1 - Present Momentum (100-150 words):**
        - Describe the current state and ongoing trends
        - Establish the "Push of the Present" foundation
        - Set the context for continuation rather than transformation
        
        **Paragraph 2 - Primary Drivers (100-140 words):**
        - Focus on the highest certainty, highest impact drivers
        - Explain how these forces reinforce current trajectories
        - Show momentum building from existing patterns
        
        **Paragraph 3 - Secondary Drivers & Evolution (100-130 words):**
        - Include additional drivers that support the baseline path
        - Show how the domain evolves within existing frameworks
        - Demonstrate gradual rather than revolutionary change
        
        **Paragraph 4 - Baseline Future State (100-120 words):**
        - Synthesize into a coherent "most likely" future
        - Emphasize continuation and extension of current trends
        - Position as the foundation before exploring alternatives
        
        **Writing Style:**
        - Narrative and story-like, but grounded in evidence
        - Confident but not overly optimistic
        - Focus on "what's most likely" rather than "what's possible"
        - Use concrete details from the domain context
        
        **Critical Focus:**
        - This is NOT about transformation or disruption
        - This IS about logical extension of current momentum
        - Emphasize high-certainty, predictable developments
        - Set up the contrast for later alternative scenarios

        **CRITICAL FORMATTING RULES:**
        - Return valid JSON with no markdown code blocks
        - The scenario_text must be a single continuous string
        - Replace all actual newlines in text with \\n escape sequences
        - Do NOT put line breaks immediately after opening quotes
        
        Format your response as JSON:
        {{
            "scenario_title": "Descriptive title for the baseline scenario",
            "timeframe": "2025-2030" or appropriate timeframe,
            "scenario_text": "Paragraph 1 content here.\\n\\nParagraph 2 content here.\\n\\nParagraph 3 content here.",
            "key_assumptions": ["assumption 1", "assumption 2", "assumption 3"],
            "dominant_drivers": ["driver 1", "driver 2", "driver 3"],
            "scenario_type": "Baseline/Continuation"
        }}
        
        Ensure the scenario text is exactly 3-4 paragraphs with natural narrative flow.
        """
        
        try:
            chat_completion = self.client.chat.completions.create(
                messages=[
                    {
                        "role": "system",
                        "content": """You are a senior strategic foresight analyst specializing in baseline scenario development. 
                        You excel at creating compelling "business-as-usual" narratives that extrapolate current trends and 
                        high-certainty drivers into plausible continuation scenarios. Your scenarios are grounded, realistic, 
                        and set the foundation for exploring alternative futures. Always respond with valid JSON."""
                    },
                    {"role": "user", "content": prompt}
                ],
                model=self.model,
                max_tokens=1500,
                temperature=0.6  # Lower temperature for more consistent baseline scenarios
            )
            
            response_text = chat_completion.choices[0].message.content
            parsed_result = self._parse_json_response(response_text)

            
            # Validate required fields
            required_fields = ['scenario_title', 'scenario_text', 'key_assumptions', 'scenario_type']
            for field in required_fields:
                if field not in parsed_result:
                    parsed_result[field] = f"Generated {field} for {domain}"
                            
            # Ensure scenario_text exists and is reasonable length
            if not parsed_result.get('scenario_text') or len(parsed_result['scenario_text']) < 100:
                parsed_result['scenario_text'] = f"Baseline scenario for {domain} continues current trends with gradual evolution driven by existing momentum and high-certainty factors."
            
            return parsed_result
                
        except Exception as e:
            return {"error": f"Failed to generate baseline scenario: {str(e)}"}
            
    def generate_driver_outcomes(self, domain: str, triangle_2_0_data: Dict, baseline_data: Dict, phase1_data: Dict = None) -> Dict[str, Any]:
        """Generate Driver Outcomes by 'bending' each driver, uncertainty, and narrative to archetypal scenarios."""
        
        # Extract elements from Futures Triangle 2.0
        drivers = triangle_2_0_data.get('drivers', [])
        uncertainties = triangle_2_0_data.get('uncertainties', [])
        narratives = triangle_2_0_data.get('narratives', [])
        
        # Project context
        project_name = phase1_data.get('project_name', domain) if phase1_data else domain
        
        # Format baseline context
        baseline_context = f"""
        BASELINE SCENARIO: {baseline_data.get('scenario_title', '')}
        Timeframe: {baseline_data.get('timeframe', '2025-2030')}
        Key Assumptions: {', '.join(baseline_data.get('key_assumptions', []))}
        """
        
        prompt = f"""
        You are a strategic foresight expert creating Driver Outcomes for "{domain}" by "bending" elements from Futures Triangle 2.0 analysis into different archetypal scenarios.

        PROJECT CONTEXT:
        Project: {project_name}
        Domain: {domain}
        
        {baseline_context}
        
        DRIVER OUTCOMES METHODOLOGY:
        Take each Driver, Uncertainty, and Narrative and "bend" them into 3 archetypal futures:
        
        1. **COLLAPSE/DECLINE** - Systems break down, failures cascade, things get worse
        2. **NEW EQUILIBRIUM** - Adaptive change, new stable patterns, reformed systems  
        3. **TRANSFORMATION** - Breakthrough innovation, paradigm shifts, fundamental change
        
        ELEMENTS TO BEND:
        
        DRIVERS (Major Forces):
        {chr(10).join([f"- {d.get('name', '')}: {d.get('description', '')} (Impact: {d.get('impact_level', '')}, Certainty: {d.get('certainty', '')})" for d in drivers[:6]])}
        
        UNCERTAINTIES (Pivot Points):
        {chr(10).join([f"- {u.get('name', '')}: {u.get('description', '')}" for u in uncertainties[:5]])}
        
        NARRATIVES (Stories):
        {chr(10).join([f"- {n.get('name', '')} ({n.get('type', '')}): {n.get('description', '')}" for n in narratives[:4]])}
        
        ARCHETYPE DEFINITIONS:
        
        **COLLAPSE/DECLINE:**
        - Systems fail, break down, or regress
        - Negative feedback loops dominate
        - Resources become scarce, trust erodes
        - Institutions lose effectiveness
        - Focus: "What goes wrong?"
        
        **NEW EQUILIBRIUM:**
        - Adaptive responses create stability
        - Systems reform and find balance
        - Gradual improvement within existing frameworks
        - Incremental innovation and adjustment
        - Focus: "How do we adapt?"
        
        **TRANSFORMATION:**
        - Breakthrough innovations emerge
        - Fundamental paradigm shifts occur
        - New systems replace old ones
        - Exponential positive change
        - Focus: "What becomes possible?"
        
        OUTCOME REQUIREMENTS:
        - Each element gets 3 outcomes (one per archetype)
        - Outcomes should be 2-3 sentences each
        - Stay grounded in the domain context
        - Show how the same force creates different futures
        - Make outcomes specific and plausible within each archetype
        
        Format as JSON:
        {{
            "driver_outcomes": [
                {{
                    "driver_id": "D1",
                    "driver_name": "Driver name from Triangle 2.0",
                    "baseline_trajectory": "How this plays out in baseline",
                    "outcomes": [
                        {{
                            "archetype": "Collapse/Decline",
                            "outcome_text": "2-3 sentence description of how this driver manifests in a collapse scenario",
                            "key_impacts": ["impact 1", "impact 2", "impact 3"]
                        }},
                        {{
                            "archetype": "New Equilibrium", 
                            "outcome_text": "2-3 sentence description of how this driver manifests in adaptive change",
                            "key_impacts": ["impact 1", "impact 2", "impact 3"]
                        }},
                        {{
                            "archetype": "Transformation",
                            "outcome_text": "2-3 sentence description of how this driver creates breakthrough change", 
                            "key_impacts": ["impact 1", "impact 2", "impact 3"]
                        }}
                    ]
                }}
            ],
            "uncertainty_outcomes": [
                {{
                    "uncertainty_id": "U1",
                    "uncertainty_name": "Uncertainty name from Triangle 2.0",
                    "key_variables": ["var1", "var2"],
                    "outcomes": [
                        {{
                            "archetype": "Collapse/Decline",
                            "outcome_text": "How this uncertainty resolves in a collapse scenario",
                            "resolution_direction": "Which way the uncertainty tips"
                        }},
                        {{
                            "archetype": "New Equilibrium",
                            "outcome_text": "How this uncertainty resolves in adaptive change",
                            "resolution_direction": "Which way the uncertainty tips"
                        }},
                        {{
                            "archetype": "Transformation", 
                            "outcome_text": "How this uncertainty resolves in transformation",
                            "resolution_direction": "Which way the uncertainty tips"
                        }}
                    ]
                }}
            ],
            "narrative_outcomes": [
                {{
                    "narrative_id": "N1",
                    "narrative_name": "Narrative name from Triangle 2.0",
                    "narrative_type": "Dominant/Emerging/Alternative",
                    "outcomes": [
                        {{
                            "archetype": "Collapse/Decline",
                            "outcome_text": "How this narrative evolves in collapse",
                            "narrative_shift": "What story dominates"
                        }},
                        {{
                            "archetype": "New Equilibrium",
                            "outcome_text": "How this narrative evolves in adaptation", 
                            "narrative_shift": "What story dominates"
                        }},
                        {{
                            "archetype": "Transformation",
                            "outcome_text": "How this narrative evolves in transformation",
                            "narrative_shift": "What story dominates"
                        }}
                    ]
                }}
            ],
            "cross_archetype_insights": {{
                "collapse_patterns": ["Common themes across collapse outcomes"],
                "equilibrium_patterns": ["Common themes across equilibrium outcomes"], 
                "transformation_patterns": ["Common themes across transformation outcomes"],
                "leverage_points": ["Key intervention points that could shift outcomes between archetypes"]
            }}
        }}
        
        CRITICAL: Ensure each element from Triangle 2.0 is "bent" to show how the SAME force creates different futures under different conditions.
        """
        
        try:
            chat_completion = self.client.chat.completions.create(
                messages=[
                    {
                        "role": "system", 
                        "content": """You are a senior strategic foresight analyst specializing in archetypal scenario development. 
                        You excel at taking identified drivers, uncertainties, and narratives and showing how they manifest 
                        differently across collapse, equilibrium, and transformation archetypes. You create plausible, 
                        specific outcomes that demonstrate how the same forces can lead to very different futures. 
                        Always respond with valid, complete JSON."""
                    },
                    {"role": "user", "content": prompt}
                ],
                model=self.model,
                max_tokens=4000,
                temperature=0.7
            )
            
            response_text = chat_completion.choices[0].message.content
            
            # Clean response - remove any markdown code blocks
            response_text = response_text.strip()
            if response_text.startswith('```json'):
                response_text = response_text.replace('```json', '').replace('```', '').strip()
            elif response_text.startswith('```'):
                response_text = response_text.replace('```', '').strip()
                
            parsed_result = self._parse_json_response(response_text)
            
            # Validate required sections exist
            required_sections = ['driver_outcomes', 'uncertainty_outcomes', 'narrative_outcomes']
            for section in required_sections:
                if section not in parsed_result:
                    parsed_result[section] = []
            
            # Ensure we have cross-archetype insights
            if 'cross_archetype_insights' not in parsed_result:
                parsed_result['cross_archetype_insights'] = {
                    'collapse_patterns': ['System failures', 'Resource scarcity', 'Trust erosion'],
                    'equilibrium_patterns': ['Adaptive responses', 'Gradual reform', 'Balanced solutions'],
                    'transformation_patterns': ['Breakthrough innovation', 'Paradigm shifts', 'Exponential change'],
                    'leverage_points': ['Policy interventions', 'Technology adoption', 'Cultural shifts']
                }
            
            return parsed_result
            
        except Exception as e:
            return {"error": f"Failed to generate driver outcomes: {str(e)}"}


#much better ----solved No more lazy “#2” titles
    def generate_alternative_scenarios(self, domain: str, selected_archetypes: Dict, baseline_data: Dict, driver_outcomes: Dict, triangle_2_0_data: Dict = None) -> Dict[str, Any]:
        """Generate alternative scenarios based on selected archetypes."""
        
        # Archetype definitions
        archetype_definitions = {
            "Collapse": "System breakdown, failures cascade, institutions lose effectiveness, negative feedback loops dominate, resources become scarce, trust erodes",
            "New Equilibrium": "Adaptive responses create stability, systems reform and find balance, gradual improvement within existing frameworks, incremental innovation and adjustment", 
            "Transformation": "Breakthrough innovations emerge, fundamental paradigm shifts occur, new systems replace old ones, exponential positive change"
        }
        
        scenarios = []
        
        for archetype, count in selected_archetypes.items():
            if count > 0:
                for i in range(count):
                    scenario = self._generate_single_scenario(
                        domain=domain,
                        archetype=archetype, 
                        archetype_definition=archetype_definitions.get(archetype, ""),
                        baseline_data=baseline_data,
                        driver_outcomes=driver_outcomes,
                        triangle_2_0_data=triangle_2_0_data,
                        scenario_number=i+1
                    )
                    scenarios.append(scenario)
        
        return {"scenarios": scenarios}

    def _generate_single_scenario(self, domain: str, archetype: str, archetype_definition: str, 
                                baseline_data: Dict, driver_outcomes: Dict, triangle_2_0_data: Dict = None, 
                                scenario_number: int = 1) -> Dict:
        """Generate a single scenario narrative with improved diversity."""
        
        # Extract key context
        baseline_text = baseline_data.get('scenario_text', '')
        baseline_title = baseline_data.get('scenario_title', '')
        
        # Get driver outcomes for this archetype
        relevant_outcomes = []
        for driver in driver_outcomes.get('driver_outcomes', []):
            for outcome in driver.get('outcomes', []):
                outcome_archetype = outcome.get('archetype', '').lower().replace(' ', '')
                target_archetype = archetype.lower().replace(' ', '')
                if target_archetype in outcome_archetype or outcome_archetype in target_archetype:
                    relevant_outcomes.append(f"• {driver.get('driver_name', '')}: {outcome.get('outcome_text', '')}")
        
        # If no relevant outcomes found, get first few driver outcomes
        if not relevant_outcomes:
            for driver in driver_outcomes.get('driver_outcomes', [])[:3]:
                for outcome in driver.get('outcomes', [])[:1]:
                    relevant_outcomes.append(f"• {driver.get('driver_name', '')}: {outcome.get('outcome_text', '')}")

        # Create scenario-specific focus areas to ensure diversity
        focus_areas = {
            "Collapse": [
                "financial system breakdown and economic collapse",
                "institutional failure and governance breakdown", 
                "technological obsolescence and infrastructure decay",
                "social fragmentation and cultural alienation"
            ],
            "New Equilibrium": [
                "sustainable development and environmental stewardship",
                "inclusive governance and democratic reforms",
                "regional cooperation and diplomatic balance",
                "tradition preservation with selective innovation"
            ],
            "Transformation": [
                "breakthrough technological revolution and digitization",
                "global democratization and grassroots expansion",
                "radical business model innovation and new economics",
                "social impact revolution and cultural transformation"
            ]
        }

        unique_drivers_per_scenario = {
            "Collapse": [
                "sponsorship withdrawal, broadcasting revenue collapse, financial mismanagement",
                "regulatory conflicts, visa restrictions, political tensions between nations",
                "aging infrastructure, resistance to new technology, equipment failures",
                "generational disconnect, competing entertainment, loss of cultural relevance"
            ],
            "New Equilibrium": [
                "carbon-neutral stadiums, renewable energy adoption, environmental regulations",
                "stakeholder representation, transparent governance, democratic decision-making",
                "cross-border partnerships, measured expansion, diplomatic cricket initiatives", 
                "heritage conservation, selective tech integration, cultural preservation"
            ],
            "Transformation": [
                "AI analytics, VR experiences, blockchain ticketing, digital fan engagement",
                "non-traditional markets, grassroots accessibility, global talent mobility",
                "subscription models, fan ownership, cryptocurrency integration, direct investment",
                "gender equality initiatives, community development, social change catalyst"
            ]
        }

        # Select focus and unique drivers based on scenario number
        focus_list = focus_areas.get(archetype, ["general system changes"])
        selected_focus = focus_list[(scenario_number - 1) % len(focus_list)]

        drivers_list = unique_drivers_per_scenario.get(archetype, ["general drivers"])
        unique_drivers = drivers_list[(scenario_number - 1) % len(drivers_list)]

        prompt = f"""Create scenario #{scenario_number} for {archetype} archetype in {domain}.

        ARCHETYPE: {archetype} - {archetype_definition}
        UNIQUE FOCUS: This scenario must focus EXCLUSIVELY on {selected_focus}
        UNIQUE DRIVERS TO EMPHASIZE: {unique_drivers}

        BASELINE CONTEXT:
        {baseline_text[:400]}

        REQUIRED DRIVER OUTCOMES TO INTEGRATE:
        {chr(10).join(relevant_outcomes[:4])}

        CRITICAL DIVERSITY REQUIREMENTS:
        - This scenario must be COMPLETELY DIFFERENT from other {archetype} scenarios
        - Focus ONLY on {selected_focus} - do not mix with other focus areas
        - Emphasize these unique drivers: {unique_drivers}
        - Create a unique storyline with different triggers, progression, and outcomes
        - The scenario title must be creative and reflect {selected_focus} (no generic numbering)
        - Probability assessment must vary and be justified
        - All factors, assumptions, and indicators must be unique to this specific focus

        STRUCTURE (EXACTLY 4 paragraphs, 300-400 words total):
        1. Initial conditions specific to {selected_focus} (2025-2026) - 75-100 words
        2. Key developments driven by {unique_drivers} (2027-2028) - 75-100 words
        3. Full manifestation of {selected_focus} (2029-2030) - 75-100 words
        4. Final system state shaped by {selected_focus} (2030) - 75-100 words

        **CRITICAL FORMATTING RULES:**
        - Return valid JSON with no markdown code blocks
        - The scenario_text must be a single continuous string
        - Replace all actual newlines in text with \\n\\n escape sequences
        - Do NOT put line breaks immediately after opening quotes
        - Keep total word count between 300-400 words
        - Each paragraph should be 75-100 words maximum

        Return ONLY valid JSON:
        {{
            "scenario_title": "Creative title reflecting {selected_focus} (no numbering)",
            "archetype": "{archetype}",
            "timeframe": "2025-2030", 
            "scenario_text": "Four paragraphs separated by \\n\\n, each 75-100 words, total 300-400 words...",
            "key_factors": ["factor specific to {selected_focus}", "factor related to {unique_drivers}", "factor 3", "factor 4", "factor 5"],
            "critical_assumptions": ["assumption about {unique_drivers}", "assumption about {selected_focus}", "third unique assumption"],
            "probability_assessment": "Low/Medium/High - [JUSTIFY WHY based on {unique_drivers}]",
            "key_indicators": ["early warning for {unique_drivers}", "indicator for {selected_focus}", "third specific indicator"]
        }}

        ENSURE: Everything must be unique to {selected_focus} and driven by {unique_drivers}. No overlap with other scenarios."""

        try:
            response = self.client.chat.completions.create(
                messages=[
                    {
                        "role": "system", 
                        "content": f"""Create a completely unique {archetype} scenario focused ONLY on {selected_focus}. 
                        Emphasize these unique drivers: {unique_drivers}. Generate a creative, metaphorical title (no numbering). 
                        Vary probability assessment (Low/Medium/High) with detailed justification. Make all factors, assumptions, 
                        and indicators scenario-specific. Always respond with valid JSON only. Never use markdown code blocks."""
                    },
                    {"role": "user", "content": prompt}
                ],
                model=self.model,
                max_tokens=1200,  # Reduced to encourage conciseness
                temperature=0.7,  # Reduced for better structure adherence
                # Remove response_format since you're handling JSON parsing manually
            )
            
            response_text = response.choices[0].message.content.strip()
            
            # CRITICAL FIX: Use your robust parser instead of json.loads()
            parsed_result = self._parse_json_response(response_text)
            
            # Check if parsing failed (empty dict or error)
            if not parsed_result or parsed_result.get('error'):
                raise ValueError("JSON parsing failed")
            
            # Validate and ensure required fields
            if not parsed_result.get('scenario_text') or len(parsed_result.get('scenario_text', '')) < 200:
                raise ValueError("Scenario text too short or missing")
            
            # Keep the AI-generated title as-is (no numbering fallback)
            title = parsed_result.get('scenario_title', f"Untitled {archetype} Scenario")
            parsed_result['scenario_title'] = title
                
            # Set proper defaults
            parsed_result.setdefault('archetype', archetype)
            parsed_result.setdefault('timeframe', '2025-2030')
            parsed_result.setdefault('key_factors', [])
            parsed_result.setdefault('critical_assumptions', [])
            parsed_result.setdefault('probability_assessment', 'Medium')
            parsed_result.setdefault('key_indicators', [])
            
            return parsed_result
            
        except Exception as e:
            print(f"Scenario generation error: {str(e)}")
            return self._generate_simple_scenario(domain, archetype, scenario_number, selected_focus)

    def _generate_simple_scenario(self, domain: str, archetype: str, scenario_number: int, focus_area: str = "") -> Dict:
        """Fallback simple scenario generation with focus area."""
        
        simple_prompt = f"""Create {archetype} scenario #{scenario_number} for {domain} (2025-2030).

    Focus on: {focus_area or archetype.lower()}

    Write 3 paragraphs showing progression over time.

    Return JSON:
    {{
        "scenario_title": "Unique title for scenario #{scenario_number}",
        "archetype": "{archetype}",
        "timeframe": "2025-2030",
        "scenario_text": "3 paragraph narrative...",
        "key_factors": ["factor1", "factor2", "factor3"],
        "critical_assumptions": ["assumption1", "assumption2"], 
        "probability_assessment": "Low/Medium/High",
        "key_indicators": ["indicator1", "indicator2"]
    }}"""
        
        try:
            response = self.client.chat.completions.create(
                messages=[
                    {"role": "system", "content": f"Create unique scenario focusing on {focus_area}. Generate creative title (no numbering). Vary probability assessment with justification."},
                    {"role": "user", "content": simple_prompt}
                ],
                model=self.model,
                max_tokens=1500,
                temperature=0.9,
                response_format={"type": "json_object"}
            )
            
            result = json.loads(response.choices[0].message.content.strip())
            
            # Ensure unique title
            title = result.get('scenario_title', f"{archetype} Scenario {scenario_number}")
            if scenario_number > 1:
                title = f"{title} #{scenario_number}"
            result['scenario_title'] = title
            
            return result
            
        except Exception as e:
            print(f"Simple scenario generation failed: {str(e)}")
            return {
                "scenario_title": f"{archetype} Focus: {focus_area} #{scenario_number}",
                "archetype": archetype,
                "timeframe": "2025-2030",
                "scenario_text": f"This {archetype.lower()} scenario explores how {domain} evolves through {focus_area} from 2025 to 2030. Early indicators emerge by 2026, with key developments unfolding through 2027-2028. By 2030, the {archetype.lower()} pattern is fully established, demonstrating the impact of {focus_area} on the system's evolution.",
                "key_factors": [f"{focus_area}", f"{domain} dynamics", "System responses"],
                "critical_assumptions": [f"{focus_area} continues as expected", "Key stakeholders adapt accordingly"],
                "probability_assessment": ["Low", "Medium", "High"][scenario_number % 3],
                "key_indicators": [f"Signs of {focus_area}", "System metric changes"]
            }

# Utility functions for Streamlit integration
def get_api_key():
    """Get Groq API key from environment or user input."""
    return os.getenv('GROQ_API_KEY', '')

def initialize_processor():
    """Initialize the DRI Foresight processor."""
    api_key = get_api_key()
    if not api_key:
        raise ValueError("GROQ_API_KEY environment variable not set")
    return DRIForesightProcessor(api_key)






